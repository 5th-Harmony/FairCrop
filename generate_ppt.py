import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

if sys.platform.startswith("win"):
    try:
        sys.stdout.reconfigure(encoding="utf-8")
        sys.stderr.reconfigure(encoding="utf-8")
    except Exception:
        pass

def create_presentation(output_path="SIH2026_AgriLink_Presentation.pptx"):
    prs = Presentation()
    # 16:9 Widescreen standard
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Theme Colors
    PRIMARY_DARK = RGBColor(15, 81, 50)      # Deep Emerald Green #0F5132
    PRIMARY_LIGHT = RGBColor(25, 135, 84)    # Green #198754
    ACCENT_GOLD = RGBColor(217, 119, 6)      # Amber/Gold #D97706
    BG_LIGHT = RGBColor(248, 250, 252)       # Light Slate #F8FAFC
    CARD_BG = RGBColor(255, 255, 255)        # Pure White
    TEXT_DARK = RGBColor(30, 41, 59)         # Slate 800
    TEXT_MUTED = RGBColor(100, 116, 139)     # Slate 500
    CARD_BORDER = RGBColor(226, 232, 240)    # Slate 200

    def set_slide_background(slide, color):
        bg_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = color
        bg_shape.line.fill.background()
        return bg_shape

    def add_header(slide, title_text, category_text="SMART INDIA HACKATHON 2026 | PROBLEM STATEMENT: MARKET LINKAGES & PRICE DISCOVERY"):
        # Top banner
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.4))
        tf_c = cat_box.text_frame
        tf_c.word_wrap = True
        p_c = tf_c.paragraphs[0]
        p_c.text = category_text.upper()
        p_c.font.size = Pt(10)
        p_c.font.bold = True
        p_c.font.color.rgb = ACCENT_GOLD

        # Title
        t_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.8))
        tf_t = t_box.text_frame
        tf_t.word_wrap = True
        p_t = tf_t.paragraphs[0]
        p_t.text = title_text
        p_t.font.size = Pt(24)
        p_t.font.bold = True
        p_t.font.color.rgb = PRIMARY_DARK

    def add_card(slide, left, top, width, height, title, content_items, accent_color=None, bg_color=CARD_BG):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
        shape.line.color.rgb = accent_color if accent_color else CARD_BORDER
        shape.line.width = Pt(1.5 if accent_color else 1.0)

        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.25)
        tf.margin_right = Inches(0.25)
        tf.margin_top = Inches(0.25)
        tf.margin_bottom = Inches(0.25)

        if title:
            p_title = tf.paragraphs[0]
            p_title.text = title
            p_title.font.size = Pt(16)
            p_title.font.bold = True
            p_title.font.color.rgb = PRIMARY_DARK if not accent_color else accent_color
            p_title.space_after = Pt(10)

        for i, item in enumerate(content_items):
            p = tf.add_paragraph() if (title or i > 0) else tf.paragraphs[0]
            p.text = item
            p.font.size = Pt(12)
            p.font.color.rgb = TEXT_DARK
            p.space_after = Pt(6)

        return shape

    # ==========================================
    # SLIDE 1: Title Slide (Dark Theme)
    # ==========================================
    s1 = prs.slides.add_slide(blank_layout)
    set_slide_background(s1, PRIMARY_DARK)

    # Sub-heading badge
    badge = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.2), Inches(5.5), Inches(0.45))
    badge.fill.solid()
    badge.fill.fore_color.rgb = ACCENT_GOLD
    badge.line.fill.background()
    tf_b = badge.text_frame
    p_b = tf_b.paragraphs[0]
    p_b.text = "SMART INDIA HACKATHON (SIH) 2026"
    p_b.font.size = Pt(12)
    p_b.font.bold = True
    p_b.font.color.rgb = RGBColor(255, 255, 255)
    p_b.alignment = PP_ALIGN.CENTER

    # Main Project Title
    t1_box = s1.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(2.2))
    tf1 = t1_box.text_frame
    tf1.word_wrap = True
    p1 = tf1.paragraphs[0]
    p1.text = "AgriLink: Transparent Price Discovery &\nMarket Linkages Platform"
    p1.font.size = Pt(36)
    p1.font.bold = True
    p1.font.color.rgb = RGBColor(255, 255, 255)
    p1.space_after = Pt(12)

    p1_sub = tf1.add_paragraph()
    p1_sub.text = "Empowering Smallholders & FPOs with AI Market Intelligence, Smart Matchmaking & Guaranteed Escrow Settlements"
    p1_sub.font.size = Pt(16)
    p1_sub.font.color.rgb = RGBColor(209, 250, 229)

    # Bottom Stat Cards
    stat1 = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(4.7), Inches(3.6), Inches(1.8))
    stat1.fill.solid()
    stat1.fill.fore_color.rgb = RGBColor(20, 95, 60)
    stat1.line.color.rgb = PRIMARY_LIGHT
    tf_s1 = stat1.text_frame
    tf_s1.word_wrap = True
    p_s1 = tf_s1.paragraphs[0]
    p_s1.text = "🌾 Farmer Realization"
    p_s1.font.size = Pt(14)
    p_s1.font.bold = True
    p_s1.font.color.rgb = ACCENT_GOLD
    p_s1_sub = tf_s1.add_paragraph()
    p_s1_sub.text = "+15% to +25%\nHigher price realization via direct institutional trade."
    p_s1_sub.font.size = Pt(12)
    p_s1_sub.font.color.rgb = RGBColor(255, 255, 255)

    stat2 = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.8), Inches(4.7), Inches(3.6), Inches(1.8))
    stat2.fill.solid()
    stat2.fill.fore_color.rgb = RGBColor(20, 95, 60)
    stat2.line.color.rgb = PRIMARY_LIGHT
    tf_s2 = stat2.text_frame
    tf_s2.word_wrap = True
    p_s2 = tf_s2.paragraphs[0]
    p_s2.text = "⚡ Real-Time Mandis"
    p_s2.font.size = Pt(14)
    p_s2.font.bold = True
    p_s2.font.color.rgb = ACCENT_GOLD
    p_s2_sub = tf_s2.add_paragraph()
    p_s2_sub.text = "Live APMC Prices +\n7-Day AI Price Forecasts with Hold/Sell recommendations."
    p_s2_sub.font.size = Pt(12)
    p_s2_sub.font.color.rgb = RGBColor(255, 255, 255)

    stat3 = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.8), Inches(4.7), Inches(3.6), Inches(1.8))
    stat3.fill.solid()
    stat3.fill.fore_color.rgb = RGBColor(20, 95, 60)
    stat3.line.color.rgb = PRIMARY_LIGHT
    tf_s3 = stat3.text_frame
    tf_s3.word_wrap = True
    p_s3 = tf_s3.paragraphs[0]
    p_s3.text = "🔒 Zero Default Risk"
    p_s3.font.size = Pt(14)
    p_s3.font.bold = True
    p_s3.font.color.rgb = ACCENT_GOLD
    p_s3_sub = tf_s3.add_paragraph()
    p_s3_sub.text = "7-Stage Escrow Lifecycle\nGuaranteed payment security & dispute resolution."
    p_s3_sub.font.size = Pt(12)
    p_s3_sub.font.color.rgb = RGBColor(255, 255, 255)

    # ==========================================
    # SLIDE 2: Problem Statement & Ground Realities
    # ==========================================
    s2 = prs.slides.add_slide(blank_layout)
    set_slide_background(s2, BG_LIGHT)
    add_header(s2, "The Core Problem: Agricultural Market Inefficiencies in India")

    add_card(s2, 0.8, 1.7, 3.6, 5.0, "🔴 Severe Information Asymmetry", [
        "• Smallholder farmers lack real-time visibility into price fluctuations across nearby Mandis and regional buyers.",
        "• Price discovery is dominated by local middlemen with opaque margins (taking 15-30% cut).",
        "• Farmers lack predictive insights to decide optimal sale windows (Hold vs. Sell)."
    ], accent_color=RGBColor(220, 38, 38))

    add_card(s2, 4.8, 1.7, 3.6, 5.0, "🔴 Distress Sales & Post-Harvest Losses", [
        "• Immediate liquidity needs force farmers into distress sales right after harvest.",
        "• Lack of integrated cold storage and transparent logistics increases spoilage and quality downgrades.",
        "• Post-harvest losses in perishables exceed 20-25% annually."
    ], accent_color=RGBColor(220, 38, 38))

    add_card(s2, 8.8, 1.7, 3.6, 5.0, "🔴 Buyer Aggregation & Trust Deficit", [
        "• Institutional buyers (processors, exporters, retail chains) struggle to aggregate consistent volumes from individual farms.",
        "• Quality parameter verification (moisture, grading) is non-standardized.",
        "• High risk of payment delays, contract defaults, and lack of grievance redressal."
    ], accent_color=RGBColor(220, 38, 38))

    # ==========================================
    # SLIDE 3: Solution Architecture - AgriLink Platform
    # ==========================================
    s3 = prs.slides.add_slide(blank_layout)
    set_slide_background(s3, BG_LIGHT)
    add_header(s3, "Our Solution: AgriLink Integrated Value Chain Platform")

    add_card(s3, 0.8, 1.7, 3.6, 2.5, "1. 🔐 Multi-Tenant Identity", [
        "• Role-based portals: Farmers, FPOs, Institutional Buyers, Logistics & Admins.",
        "• Auto-verification for farmers; GSTIN verification for corporate buyers.",
        "• Dual login support: Phone Number / Email."
    ])

    add_card(s3, 4.8, 1.7, 3.6, 2.5, "2. 🌾 Lot & Quality Grading", [
        "• Digital lot creation with standardized grades (Grade A/B/C, Organic, Premium).",
        "• Moisture percentage, harvest dates & storage parameters.",
        "• Batch photo proof and geo-location tracking."
    ])

    add_card(s3, 8.8, 1.7, 3.6, 2.5, "3. 🎯 Smart Matchmaking Engine", [
        "• Ranks listings for buyers using multi-parameter compatibility scoring.",
        "• Evaluates crop type, volume fulfilment, budget savings, quality & proximity.",
        "• Bidding, counter-offers, and real-time negotiations."
    ])

    add_card(s3, 0.8, 4.5, 3.6, 2.5, "4. 💳 7-Stage Escrow Lifecycle", [
        "• Contract locked upon bid acceptance.",
        "• Buyer deposits funds into Escrow before dispatch.",
        "• Payout released to farmer upon quality approval."
    ])

    add_card(s3, 4.8, 4.5, 3.6, 2.5, "5. 📊 AI Market Intelligence", [
        "• Aggregated live APMC Mandi price benchmarks.",
        "• 7-Day predictive price forecast line charts.",
        "• AI-driven 'HOLD / SELL' decision advisory with reasoning."
    ])

    add_card(s3, 8.8, 4.5, 3.6, 2.5, "6. ⚖️ Dispute & Grievance Redressal", [
        "• Structured grievance ticketing with photo evidence uploads.",
        "• Escrow auto-pauses when a dispute is opened.",
        "• Fast-track administrative mediation and resolution."
    ])

    # ==========================================
    # SLIDE 4: Database Architecture & Technical Stack
    # ==========================================
    s4 = prs.slides.add_slide(blank_layout)
    set_slide_background(s4, BG_LIGHT)
    add_header(s4, "Backend Architecture & Robust Relational Schema")

    add_card(s4, 0.8, 1.7, 5.6, 5.2, "🗄️ Relational Database Model (SQLAlchemy 2.0)", [
        "• Users: Identity, multi-tenancy, credentials, KYC, geo-hierarchy (State -> District -> Tehsil -> Village).",
        "• ProduceLots: Crop name, variety, quantity_kg, expected price, Grade, moisture %, storage & harvest dates.",
        "• Offers: Bidding engine with counter-offers, expiry timers, and automated state transitions.",
        "• EscrowTransactions: Audit-trailed payment reference, delivery addresses, dispatch milestones, and statuses.",
        "• MandiPrices & PriceForecasts: Historical APMC arrival data + 7-day predictive AI time-series models.",
        "• DisputeTickets: Grievance classification, evidence attachments, and administrative resolution workflows."
    ], accent_color=PRIMARY_DARK)

    add_card(s4, 6.8, 1.7, 5.6, 5.2, "⚙️ Modern Technology Stack", [
        "• Framework: FastAPI (High-performance async Python backend with automated OpenAPI docs).",
        "• ORM Layer: SQLAlchemy 2.0 with complete relational constraints, indexes, and cascades.",
        "• Serialization & Validation: Pydantic v2 schemas for strict type safety and zero runtime leakage.",
        "• Security: Bcrypt password hashing + JWT Bearer token authentication.",
        "• Database Flexibility: Native SQLite for local zero-config testing; PostgreSQL ready for cloud production.",
        "• Seed Architecture: Realistic Indian agricultural dataset covering Nashik Onions, Khanna Wheat, Azadpur Tomatoes, and APMC benchmarks."
    ], accent_color=PRIMARY_LIGHT)

    # ==========================================
    # SLIDE 5: Smart Matchmaking & Price Discovery
    # ==========================================
    s5 = prs.slides.add_slide(blank_layout)
    set_slide_background(s5, BG_LIGHT)
    add_header(s5, "Algorithmic Smart Matchmaking & Price Discovery")

    add_card(s5, 0.8, 1.7, 5.6, 5.2, "🧠 5-Factor Match Compatibility Scoring", [
        "1. Exact Crop Relevance (30% weight): Matches buyer requirements with verified farmer listings.",
        "2. Volume Fulfillment (25% weight): Scores capacity to fulfill required batch sizes (or aggregates FPO lots).",
        "3. Price Competitiveness (20% weight): Rewards competitive pricing below buyer's maximum ceiling budget.",
        "4. Quality & Grade Matching (15% weight): Matches Grade A / Organic / Premium specifications.",
        "5. Geo-Proximity (10% weight): Minimizes transit distances, reducing transit time and carbon footprint."
    ], accent_color=ACCENT_GOLD)

    add_card(s5, 6.8, 1.7, 5.6, 5.2, "📈 AI Price Forecast & Hold/Sell Advisory", [
        "• Live Mandi Aggregation: Continuous polling and normalization of APMC modal prices.",
        "• 7-Day Trend Analysis: Machine learning forecast model projecting 7-day price trajectory with confidence intervals.",
        "• Actionable Recommendations:",
        "  - 'HOLD Advisory': When supply shortages in neighboring states indicate upcoming price spikes (+10-15%).",
        "  - 'SELL Advisory': When massive polyhouse arrivals are forecasted to crash prices within 48-72 hours.",
        "• Transparent Rationale: Delivers plain-language explanations so farmers understand why."
    ], accent_color=PRIMARY_LIGHT)

    # ==========================================
    # SLIDE 6: Escrow State Machine & Trust Lifecycle
    # ==========================================
    s6 = prs.slides.add_slide(blank_layout)
    set_slide_background(s6, BG_LIGHT)
    add_header(s6, "Trust & Guaranteed Payment: 7-Stage Escrow State Machine")

    steps = [
        ("1. INITIATED", "Bid accepted by farmer/FPO.\nContract automatically locked."),
        ("2. ESCROW DEPOSITED", "Buyer locks 100% payment\ninto secure escrow account."),
        ("3. DISPATCHED", "Produce packed & dispatched\nwith logistics tracking number."),
        ("4. DELIVERED", "Goods arrive at buyer warehouse\nfor quality inspection."),
        ("5. ESCROW RELEASED", "Buyer approves quality;\nfunds instantly sent to farmer."),
        ("6. DISPUTE / PAUSE", "Paused if grievance raised;\nadmin mediates resolution."),
        ("7. REVERT / CANCEL", "If cancelled, lot reverts\nback to AVAILABLE status.")
    ]

    for idx, (st_title, st_desc) in enumerate(steps[:4]):
        add_card(s6, 0.8 + (idx * 2.95), 1.7, 2.75, 2.4, st_title, [st_desc], accent_color=PRIMARY_DARK)

    for idx, (st_title, st_desc) in enumerate(steps[4:]):
        add_card(s6, 0.8 + (idx * 3.9), 4.4, 3.7, 2.4, st_title, [st_desc], accent_color=PRIMARY_LIGHT)

    # ==========================================
    # SLIDE 7: Impact & Measurable Outcomes
    # ==========================================
    s7 = prs.slides.add_slide(blank_layout)
    set_slide_background(s7, BG_LIGHT)
    add_header(s7, "Measurable Impact & SIH 2026 Expected Outcomes")

    add_card(s7, 0.8, 1.7, 3.6, 5.0, "🌾 Farmer Empowerment", [
        "• Improved Price Realisation: Direct institutional linkages eliminate 15-25% middleman brokerage.",
        "• Informed Sale Timing: AI advisories help farmers avoid market glut troughs.",
        "• FPO Aggregation: Enables smallholders to pool harvest into bulk institutional grade lots."
    ], accent_color=PRIMARY_DARK)

    add_card(s7, 4.8, 1.7, 3.6, 5.0, "🏭 Buyer Efficiency", [
        "• Reliable Sourcing: Guaranteed volume aggregation with verified quality specs (moisture, grade).",
        "• Transparent Logistics: End-to-end milestone tracking from farm gate to factory floor.",
        "• Reduced Procurement Cost: Direct digital offer negotiation with transparent pricing."
    ], accent_color=PRIMARY_LIGHT)

    add_card(s7, 8.8, 1.7, 3.6, 5.0, "🇮🇳 National Agricultural Growth", [
        "• Reduced Post-Harvest Spoilage: Faster market linkages cut dwell times in transit.",
        "• Complete Transaction Transparency: Zero default risk through smart escrow contracts.",
        "• Digital Ag-Stack Ready: Compliant with ONDC and Digital Public Infrastructure for Agriculture."
    ], accent_color=ACCENT_GOLD)

    # ==========================================
    # SLIDE 8: Roadmap & Next Phases (Dark Theme Conclusion)
    # ==========================================
    s8 = prs.slides.add_slide(blank_layout)
    set_slide_background(s8, PRIMARY_DARK)

    t8_box = s8.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(11.7), Inches(0.8))
    tf8 = t8_box.text_frame
    p8 = tf8.paragraphs[0]
    p8.text = "Implementation Roadmap & Future Vision"
    p8.font.size = Pt(28)
    p8.font.bold = True
    p8.font.color.rgb = RGBColor(255, 255, 255)

    add_card(s8, 0.8, 1.8, 3.6, 4.8, "Phase 1: Database & Core (Done ✅)", [
        "• Full SQLAlchemy 2.0 ORM models.",
        "• Multi-tenant role system with auto-KYC.",
        "• Pydantic v2 schemas & database seeder.",
        "• 100% integrity validation tests passed."
    ], bg_color=RGBColor(20, 95, 60))

    add_card(s8, 4.8, 1.8, 3.6, 4.8, "Phase 2: APIs & Intelligence (Next 🚀)", [
        "• FastAPI CRUD & negotiation routes.",
        "• Real-time WebSocket bid notifications.",
        "• ML price prediction inference engine.",
        "• Escrow payment gateway integration."
    ], bg_color=RGBColor(20, 95, 60))

    add_card(s8, 8.8, 1.8, 3.6, 4.8, "Phase 3: Scale & ONDC Integration", [
        "• Multilingual voice bots (Hindi, Marathi, Punjabi).",
        "• ONDC Agri-protocol buyer/seller discovery.",
        "• IoT & Satellite moisture/yield assessment.",
        "• Kisan Credit Card (KCC) micro-financing."
    ], bg_color=RGBColor(20, 95, 60))

    # Save presentation
    prs.save(output_path)
    print(f"🎉 PowerPoint presentation created successfully at: {output_path}")

if __name__ == "__main__":
    create_presentation()
