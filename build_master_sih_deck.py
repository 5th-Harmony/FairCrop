# -*- coding: utf-8 -*-
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_master_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Theme Colors
    NAVY_TITLE = RGBColor(16, 44, 87)       # #102C57
    EMERALD = RGBColor(15, 81, 50)          # #0F5132
    FOOTER_BLUE = RGBColor(30, 115, 190)    # #1E73BE
    BG_LIGHT = RGBColor(248, 250, 252)       # #F8FAFC
    CARD_BG = RGBColor(255, 255, 255)
    TEXT_DARK = RGBColor(33, 37, 41)
    TEXT_MUTED = RGBColor(100, 116, 139)
    ACCENT_ORANGE = RGBColor(255, 107, 0)   # #FF6B00
    BORDER_BLUE = RGBColor(186, 218, 245)
    BORDER_LIGHT = RGBColor(226, 232, 240)

    def set_bg(slide, color=BG_LIGHT):
        s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        s.fill.solid()
        s.fill.fore_color.rgb = color
        s.line.fill.background()
        return s

    def add_sih_header(slide, title, slide_num, total_slides=9):
        # Team Badge Top Left
        if slide_num > 1:
            team_oval = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.4), Inches(0.2), Inches(1.3), Inches(0.95))
            team_oval.fill.solid()
            team_oval.fill.fore_color.rgb = RGBColor(255, 255, 255)
            team_oval.line.color.rgb = RGBColor(50, 50, 50)
            team_oval.line.width = Pt(1.2)
            tf_team = team_oval.text_frame
            tf_team.word_wrap = True
            p1 = tf_team.paragraphs[0]
            p1.text = "TEAM"
            p1.font.size = Pt(10.5)
            p1.font.bold = True
            p1.font.color.rgb = RGBColor(20, 20, 20)
            p1.alignment = PP_ALIGN.CENTER
            p2 = tf_team.add_paragraph()
            p2.text = "METERE"
            p2.font.size = Pt(10.5)
            p2.font.bold = True
            p2.font.color.rgb = RGBColor(20, 20, 20)
            p2.alignment = PP_ALIGN.CENTER

        # Title Center
        t_box = slide.shapes.add_textbox(Inches(1.8), Inches(0.25), Inches(8.8), Inches(0.9))
        tf = t_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(26)
        p.font.bold = True
        p.font.color.rgb = NAVY_TITLE
        p.alignment = PP_ALIGN.CENTER

        # SIH Logo Top Right
        logo_box = slide.shapes.add_textbox(Inches(10.8), Inches(0.15), Inches(2.2), Inches(0.9))
        tf_l = logo_box.text_frame
        p_l = tf_l.paragraphs[0]
        p_l.text = "SMART INDIA\nHACKATHON 2026"
        p_l.font.size = Pt(11)
        p_l.font.bold = True
        p_l.font.color.rgb = NAVY_TITLE
        p_l.alignment = PP_ALIGN.CENTER

        # Footer Strip
        f_strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.05), Inches(13.333), Inches(0.45))
        f_strip.fill.solid()
        f_strip.fill.fore_color.rgb = FOOTER_BLUE
        f_strip.line.fill.background()
        tf_f = f_strip.text_frame
        p_f = tf_f.paragraphs[0]
        p_f.text = f"@SIH Idea submission- Template                                                                                                                                           {slide_num}"
        p_f.font.size = Pt(10.5)
        p_f.font.color.rgb = RGBColor(255, 255, 255)
        p_f.font.bold = True

    def add_card(slide, left, top, width, height, title, content_items, accent=BORDER_LIGHT, bg=CARD_BG, title_color=NAVY_TITLE):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        card.fill.solid()
        card.fill.fore_color.rgb = bg
        card.line.color.rgb = accent
        card.line.width = Pt(1.5)

        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.2)
        tf.margin_right = Inches(0.2)
        tf.margin_top = Inches(0.18)
        tf.margin_bottom = Inches(0.18)

        if title:
            p0 = tf.paragraphs[0]
            p0.text = title
            p0.font.size = Pt(14)
            p0.font.bold = True
            p0.font.color.rgb = title_color
            p0.space_after = Pt(6)

        for i, item in enumerate(content_items):
            p = tf.add_paragraph() if (title or i > 0) else tf.paragraphs[0]
            p.text = item
            p.font.size = Pt(10.5)
            p.font.color.rgb = TEXT_DARK
            p.space_after = Pt(4)

        return card

    # =========================================================================
    # SLIDE 1: TITLE PAGE (Official Template)
    # =========================================================================
    s1 = prs.slides.add_slide(blank_layout)
    set_bg(s1, BG_LIGHT)

    # Top Heading
    top_box = s1.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.6))
    tf1 = top_box.text_frame
    p = tf1.paragraphs[0]
    p.text = "SMART INDIA HACKATHON 2026"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = NAVY_TITLE
    p.alignment = PP_ALIGN.CENTER

    sub_t = s1.shapes.add_textbox(Inches(0.8), Inches(1.1), Inches(11.7), Inches(0.5))
    tf_sub = sub_t.text_frame
    p_sub = tf_sub.paragraphs[0]
    p_sub.text = "TITLE PAGE"
    p_sub.font.size = Pt(22)
    p_sub.font.bold = True
    p_sub.font.color.rgb = NAVY_TITLE
    p_sub.alignment = PP_ALIGN.CENTER

    # Left Metadata Box
    meta_box = s1.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(7.5), Inches(4.8))
    tf_m = meta_box.text_frame
    tf_m.word_wrap = True

    fields = [
        ("Problem Statement ID:", "SIH-1693"),
        ("Problem Statement Title:", "Strengthening market linkages and price discovery for farmers"),
        ("Theme:", "Agriculture, FoodTech & Rural Development"),
        ("PS Category:", "Software"),
        ("Team ID:", "SIH2026-METERE"),
        ("Team Name (Registered):", "METERE"),
        ("Project Name:", "FairCrop — Intelligent Agricultural Marketplace")
    ]
    for label, val in fields:
        p_lbl = tf_m.add_paragraph() if tf_m.paragraphs[0].text else tf_m.paragraphs[0]
        p_lbl.text = f"• {label}  {val}"
        p_lbl.font.size = Pt(13)
        p_lbl.font.bold = True if "Title" in label or "Name" in label else False
        p_lbl.font.color.rgb = TEXT_DARK
        p_lbl.space_after = Pt(10)

    # Right Logo / Graphic Box
    r_card = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.5), Inches(1.9), Inches(4.0), Inches(4.5))
    r_card.fill.solid()
    r_card.fill.fore_color.rgb = CARD_BG
    r_card.line.color.rgb = BORDER_BLUE
    r_card.line.width = Pt(1.5)
    tf_r = r_card.text_frame
    tf_r.word_wrap = True
    tf_r.margin_left = Inches(0.3)
    tf_r.margin_right = Inches(0.3)
    tf_r.margin_top = Inches(0.4)

    pr = tf_r.paragraphs[0]
    pr.text = "🌱 FairCrop.in"
    pr.font.size = Pt(22)
    pr.font.bold = True
    pr.font.color.rgb = EMERALD
    pr.alignment = PP_ALIGN.CENTER
    pr.space_after = Pt(14)

    pr2 = tf_r.add_paragraph()
    pr2.text = "Key Innovation:\n• Direct Farm-Gate Procurement\n• AI 7-Day Hold/Sell Price Forecast\n• 7-Stage Digital Escrow Payment\n• 22 Official Indian Languages"
    pr2.font.size = Pt(11.5)
    pr2.font.color.rgb = TEXT_DARK
    pr2.space_after = Pt(14)

    # =========================================================================
    # SLIDE 2: IDEA TITLE & SOLUTION OVERVIEW
    # =========================================================================
    s2 = prs.slides.add_slide(blank_layout)
    set_bg(s2, BG_LIGHT)
    add_sih_header(s2, "IDEA TITLE", 2)

    # Proposed Solution Banner
    banner = s2.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(11.7), Inches(0.7))
    tf_b = banner.text_frame
    tf_b.word_wrap = True
    pb = tf_b.paragraphs[0]
    pb.text = "• Proposed solution:\n“Faircrop is an intelligent agricultural marketplace for strengthening farmer-buyer linkages through transparent price information, market comparison, and demand-based recommendations.”"
    pb.font.size = Pt(12)
    pb.font.bold = True
    pb.font.color.rgb = NAVY_TITLE

    # 3 Columns
    add_card(s2, 0.8, 2.0, 3.6, 4.8, "Addressing the Problem", [
        "• Farmers lack real-time price visibility across regional mandis, limiting bargaining power.",
        "• Dominated by local middlemen who extract 15–30% hidden brokerage margins.",
        "• Inadequate storage and logistics force distress sales, causing 20–25% perishable spoilage.",
        "• Farmers face severe payment default risks from untrusted local buyers."
    ], accent=ACCENT_ORANGE)

    add_card(s2, 4.8, 2.0, 3.6, 4.8, "FairCrop Central Platform", [
        "🌐 Institutional Buyer Network:\nDirect procurement from FMCG, food processors, exporters, and retail chains.",
        "🤖 AI Market Intelligence Hub:\nAggregates 1,361+ e-NAM & Agmarknet mandis to generate real-time price forecasts.",
        "🌾 Farmer & FPO Data Inputs:\nCrop variety, quantity (quintals), moisture %, and Agmarknet Grade A/B/C specifications."
    ], accent=EMERALD)

    add_card(s2, 8.8, 2.0, 3.6, 4.8, "Solution Overview", [
        "• Demand-Based Matchmaking: Matches farmers with verified buyers based on crop, quantity, and grade.",
        "• Predictive Price Insights: AI-driven 'HOLD or SELL' advisories prevent selling in glut troughs.",
        "• 100% Escrow Security: Locks buyer funds upfront; releases payment instantly upon quality approval.",
        "• Inclusive Multilingual UI: Voice search and 22 Indian languages empower rural smallholders."
    ], accent=BORDER_BLUE)

    # =========================================================================
    # SLIDE 3: TECHNICAL APPROACH & SOFTWARE STACK
    # =========================================================================
    s3 = prs.slides.add_slide(blank_layout)
    set_bg(s3, BG_LIGHT)
    add_sih_header(s3, "TECHNICAL APPROACH", 3)

    add_card(s3, 0.8, 1.4, 5.6, 5.4, "Technologies to be Used:", [
        "• Backend Framework: Python FastAPI — ultra-fast asynchronous REST API architecture with auto-generated Swagger docs.",
        "• ORM & Database Layer: SQLAlchemy 2.0 with PostgreSQL (Production) and SQLite (Local Testing).",
        "• Strict Data Validation: Pydantic v2 schemas ensuring robust input handling and type safety.",
        "• Security & Authentication: JWT Bearer tokens + Bcrypt password hashing + RBAC (Role-Based Access Control).",
        "• AI/ML Engine: Hybrid ARIMA + LSTM time-series models for 7-day commodity price forecasting.",
        "• Live Data Feeds: Government APIs — e-NAM (1,361 unified mandis) & Agmarknet daily price feeds.",
        "• Frontend Clients: Next.js/React web dashboard, Vanilla JS client, and React Native mobile app."
    ], accent=EMERALD)

    add_card(s3, 6.8, 1.4, 5.7, 5.4, "System Architecture & Flow:", [
        "1. USER REGISTRATION: Mobile OTP + Aadhaar e-KYC (Farmer, FPO, Buyer).",
        "2. DIGITAL LOT CREATION: Farmer inputs crop, quantity, moisture %, grade.",
        "3. AI PRICE INTELLIGENCE: Fetches live mandi benchmarks + 7-day forecast.",
        "4. SMART MATCHMAKING: Matches lot with verified institutional buyers.",
        "5. DIGITAL BIDDING: Real-time negotiation and contract locking.",
        "6. 100% ESCROW LOCK: Buyer deposits funds before produce dispatch.",
        "7. TRACKED DISPATCH: GPS & cold-chain logistics milestone updates.",
        "8. DELIVERY & QUALITY CHECK: Buyer inspects produce against assay specs.",
        "9. ESCROW RELEASE: Instant payout direct to farmer bank account."
    ], accent=NAVY_TITLE)

    # =========================================================================
    # SLIDE 4: FEASIBILITY AND VIABILITY
    # =========================================================================
    s4 = prs.slides.add_slide(blank_layout)
    set_bg(s4, BG_LIGHT)
    add_sih_header(s4, "FEASIBILITY AND VIABILITY", 4)

    add_card(s4, 0.8, 1.4, 3.6, 5.4, "Feasibility Analysis", [
        "• Technical Feasibility: Built on production-proven FastAPI + PostgreSQL stack with 100% test coverage.",
        "• Operational Viability: Solves real pain points by eliminating 15–30% middleman margins and cutting transit waste.",
        "• Legal & Regulatory Fit: Aligned with the Indian Contract Act (1872) and Agmark Quality Grading Standards (1937)."
    ], accent=EMERALD)

    add_card(s4, 4.8, 1.4, 3.6, 5.4, "Challenges and Risks", [
        "• Adoption Friction: Rural farmers may have initial hesitation with digital smartphone apps.",
        "• Logistics Vulnerability: Risk of transit delays and quality degradation for perishable crops.",
        "• Price Volatility: Rapid APMC price swings could cause bid renegotiation friction."
    ], accent=ACCENT_ORANGE)

    add_card(s4, 8.8, 1.4, 3.6, 5.4, "Mitigation Strategies", [
        "• FPO Aggregation Model: Local FPOs assist village smallholders with bulk lot creation and assaying.",
        "• 7-Stage Escrow Guarantee: 100% buyer funds locked upfront; dispute mediation protocol prevents fraud.",
        "• ONDC Integration Roadmap: Plugging into the Open Network for Digital Commerce for nationwide scale."
    ], accent=BORDER_BLUE)

    # =========================================================================
    # SLIDE 5: IMPACT AND BENEFITS
    # =========================================================================
    s5 = prs.slides.add_slide(blank_layout)
    set_bg(s5, BG_LIGHT)
    add_sih_header(s5, "IMPACT AND BENEFITS", 5)

    add_card(s5, 0.8, 1.4, 3.6, 5.4, "🌾 Social Benefits", [
        "• Farmer Empowerment: Increases farm-gate price realization by 15–25% by cutting middleman cuts.",
        "• Smallholder Inclusivity: FPO collective selling gives small farmers the bargaining power of large producers.",
        "• Digital Trust: Verified buyer badges and transparent grievance handling build long-term confidence."
    ], accent=EMERALD)

    add_card(s5, 4.8, 1.4, 3.6, 5.4, "📈 Economic Benefits", [
        "• Reduced Procurement Costs: Buyers save 10–15% by procuring directly at source with verified specs.",
        "• Zero Payment Default: 100% escrow guarantee eliminates delayed payments and bad debts.",
        "• Better Price Timing: AI hold/sell signals help farmers avoid selling during market glut crashes."
    ], accent=NAVY_TITLE)

    add_card(s5, 8.8, 1.4, 3.6, 5.4, "🌍 Environmental Benefits", [
        "• 20–25% Less Food Spoilage: Direct logistics routes eliminate multi-hop mandi handling and delays.",
        "• Optimized Transportation: Consolidated FPO dispatches reduce empty freight truck runs and carbon emissions.",
        "• Efficient Resource Use: Transparent demand forecasting prevents overproduction and harvest dumping."
    ], accent=ACCENT_ORANGE)

    # =========================================================================
    # SLIDE 6: TECHNICAL VIVA CHEATSHEET (Part 1 — Architecture & AI)
    # =========================================================================
    s6 = prs.slides.add_slide(blank_layout)
    set_bg(s6, BG_LIGHT)
    add_sih_header(s6, "TECHNICAL VIVA & JUDGES Q&A (Part 1)", 6)

    add_card(s6, 0.8, 1.4, 5.6, 2.5, "Q1: Why FastAPI over Django or Flask?", [
        "• Asynchronous ASGI Engine: Handles high concurrent user traffic with ultra-low latency.",
        "• Automatic Validation: Pydantic v2 models validate incoming lot data strictly.",
        "• Swagger / OpenAPI: Auto-generates interactive API documentation out of the box."
    ], accent=BORDER_BLUE)

    add_card(s6, 6.8, 1.4, 5.7, 2.5, "Q2: Where does live mandi data come from?", [
        "• e-NAM (National Agriculture Market): 1,361+ unified mandis data across 23 states.",
        "• Agmarknet (DMI, GoI): Daily wholesale arrivals, min/max/modal prices dataset.",
        "• Open Government Data (data.gov.in): Real-time daily price updates and historical stats."
    ], accent=EMERALD)

    add_card(s6, 0.8, 4.2, 5.6, 2.6, "Q3: How does the AI Price Forecast work?", [
        "• Time-Series Predictive Models: Hybrid ARIMA + LSTM deep learning architectures.",
        "• 7-Day Price Band Forecast: Analyzes seasonality, arrival volumes, and historical patterns.",
        "• Hold vs Sell Signal: Advises farmers whether holding produce will yield higher returns."
    ], accent=ACCENT_ORANGE)

    add_card(s6, 6.8, 4.2, 5.7, 2.6, "Q4: How does the 7-Stage Escrow prevent fraud?", [
        "• 100% Upfront Deposit: Buyer funds are locked before farmer dispatches harvest.",
        "• Verified Quality Release: Funds released to farmer only after delivery inspection.",
        "• Admin Mediation: Built-in dispute resolution if moisture or grade variance occurs."
    ], accent=NAVY_TITLE)

    # =========================================================================
    # SLIDE 7: TECHNICAL VIVA CHEATSHEET (Part 2 — Rural Scale & Business)
    # =========================================================================
    s7 = prs.slides.add_slide(blank_layout)
    set_bg(s7, BG_LIGHT)
    add_sih_header(s7, "TECHNICAL VIVA & JUDGES Q&A (Part 2)", 7)

    add_card(s7, 0.8, 1.4, 5.6, 2.5, "Q5: How will rural / uneducated farmers use this?", [
        "• 22 Official Indian Languages: Native scripts (Hindi, Bengali, Telugu, Marathi, Tamil, etc.).",
        "• Voice Search: Regional dialect voice input for seamless crop and mandi lookups.",
        "• FPO Assisted Mode: FPO coordinators create lots on behalf of village members."
    ], accent=EMERALD)

    add_card(s7, 6.8, 1.4, 5.7, 2.5, "Q6: What if village internet connectivity is poor?", [
        "• Lightweight Architecture: Vanilla JS frontend without heavy framework bundle overhead.",
        "• Offline Caching: LocalStorage stores recent mandi rates and offline drafts.",
        "• Fast Loading: Optimized for low-bandwidth 2G/3G connectivity."
    ], accent=BORDER_BLUE)

    add_card(s7, 0.8, 4.2, 5.6, 2.6, "Q7: What is the Business / Revenue Model?", [
        "• 100% Free for Farmers & FPOs: Zero commission charged to smallholders.",
        "• 0.5%–1.5% Buyer Convenience Fee: Charged to large institutional bulk buyers for quality-verified sourcing.",
        "• Value-Added Services: Premium freight coordination and enterprise analytics subscriptions."
    ], accent=ACCENT_ORANGE)

    add_card(s7, 6.8, 4.2, 5.7, 2.6, "Q8: How is FairCrop different from e-NAM?", [
        "• Direct Farm-Gate Linkage: Bypasses physical mandi transport costs and auction queues.",
        "• AI Hold/Sell Intelligence: Guides farmers on optimal selling dates, unlike plain raw rate boards.",
        "• Guaranteed Escrow Security: End-to-end payment locking and direct digital settlement."
    ], accent=NAVY_TITLE)

    # =========================================================================
    # SLIDE 8: KEY NUMBERS & METRICS TO QUOTE
    # =========================================================================
    s8 = prs.slides.add_slide(blank_layout)
    set_bg(s8, BG_LIGHT)
    add_sih_header(s8, "KEY METRICS & COMPETITIVE ADVANTAGE", 8)

    add_card(s8, 0.8, 1.4, 3.6, 5.4, "📊 Key Target Metrics", [
        "• +15% to +25% Farmer Incomes: Eliminates middleman commissions and broker fees.",
        "• -20% to -25% Post-Harvest Loss: Faster direct dispatch prevents warehouse decay.",
        "• 1,361+ Mandis Integrated: Real-time price transparency across 23 states/UTs.",
        "• 100% Safe Payments: Zero default risk via upfront escrow locking."
    ], accent=EMERALD)

    add_card(s8, 4.8, 1.4, 3.6, 5.4, "⚡ Feature Comparison", [
        "Traditional Mandis vs FairCrop:",
        "❌ 15–30% Middleman Cut  ➔  ✅ 0% Farmer Commission",
        "❌ Opaque Physical Bidding  ➔  ✅ Transparent e-Bids",
        "❌ Payment Delay Risk  ➔  ✅ 100% Escrow Guarantee",
        "❌ Guesswork on Timing  ➔  ✅ AI Hold/Sell Forecasts",
        "❌ English-Only Tools  ➔  ✅ 22 Indian Languages"
    ], accent=NAVY_TITLE)

    add_card(s8, 8.8, 1.4, 3.6, 5.4, "🚀 Scalability & Roadmap", [
        "• Phase 1 (Complete ✅): FastAPI backend, ORM models, Escrow contracts, 22-language translation.",
        "• Phase 2 (Current 🚀): Live e-NAM sync, ML price forecast, Web & Mobile interfaces.",
        "• Phase 3 (Scale 🌟): ONDC Agri-Protocol integration, IoT moisture sensors, KCC credit linkage."
    ], accent=ACCENT_ORANGE)

    # =========================================================================
    # SLIDE 9: RESEARCH AND REFERENCES (Official Template)
    # =========================================================================
    s9 = prs.slides.add_slide(blank_layout)
    set_bg(s9, BG_LIGHT)
    add_sih_header(s9, "RESEARCH AND REFERENCES", 9)

    lead_box = s9.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(11.7), Inches(0.4))
    tf_lead = lead_box.text_frame
    p_lead = tf_lead.paragraphs[0]
    p_lead.text = "• Details / Links of the reference and research work foundational to FairCrop:"
    p_lead.font.size = Pt(13)
    p_lead.font.bold = True
    p_lead.font.color.rgb = NAVY_TITLE

    add_card(s9, 0.8, 1.7, 3.6, 5.1, "🏛️ Govt. Datasets & APIs", [
        "• e-NAM (National Agriculture Market): SFAC, Ministry of Agriculture. Integration standards for 1,361+ mandis.\n🔗 enam.gov.in",
        "• Agmarknet Portal (DMI, GoI): Daily wholesale arrivals, min, max, and modal price datasets.\n🔗 agmarknet.gov.in",
        "• Open Govt Data (data.gov.in): Real-time APMC price feeds and commodity statistics.\n🔗 data.gov.in",
        "• PM-KISAN & PMFBY Portals: DPI guidelines for DBT transfers & crop insurance.\n🔗 pmkisan.gov.in | pmfby.gov.in"
    ], accent=BORDER_BLUE)

    add_card(s9, 4.8, 1.7, 3.6, 5.1, "📊 Policy Reports & Studies", [
        "• NITI Aayog & NABARD (2024): “Promoting FPOs & Strengthening Digital Public Infrastructure in Indian Agriculture.” (FPOs increase farmer margins by 15–25%).",
        "• ICRIER Research Study: “Agricultural Market Integration & Price Volatility in India.” (Quantifies 18–30% middleman margin leakages across regional APMCs).",
        "• FAO & World Bank Studies: “Mitigating Post-Harvest Food Loss in Perishables via Direct Linkages.” (Direct contracting cuts perishable waste by 20–25%)."
    ], accent=EMERALD)

    add_card(s9, 8.8, 1.7, 3.6, 5.1, "⚙️ Standards & Frameworks", [
        "• ONDC Agri Protocol Specifications: Open Network for Digital Commerce specifications for unbundled agri discovery.\n🔗 ondc.org",
        "• Agmark Quality Grading Standards: Agricultural Produce Act, 1937 (Moisture %, Purity Grade A/B/C).",
        "• ML Price Forecasting Research: Time-series predictive models (ARIMA + LSTM) for APMC price trends.",
        "• Indian Contract Act (1872) & Escrow: Legal framework for digital contract lock & milestone fund release."
    ], accent=ACCENT_ORANGE)

    # Save to both file names
    prs.save("SIH2026_METERE_FairCrop.pptx")
    prs.save("SIH2026_Technical_Viva_Guide.pptx")
    prs.save("SIH2026_AgriLink_Presentation.pptx")
    print("SUCCESS: Created 9-slide Master SIH 2026 presentation with all technical, viva, and reference sections!")

if __name__ == "__main__":
    create_master_presentation()
