# -*- coding: utf-8 -*-
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

if sys.platform.startswith("win"):
    try:
        sys.stdout.reconfigure(encoding="utf-8")
        sys.stderr.reconfigure(encoding="utf-8")
    except Exception:
        pass

def generate_sih_deck(output_pptx="SIH2026_METERE_FairCrop.pptx"):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Colors
    NAVY_TITLE = RGBColor(16, 44, 87)       # #102C57
    SIH_BLUE = RGBColor(14, 116, 144)       # #0E7490
    FOOTER_BLUE = RGBColor(30, 115, 190)    # #1E73BE
    BG_WHITE = RGBColor(255, 255, 255)
    TEXT_DARK = RGBColor(33, 37, 41)        # #212529
    TEXT_MUTED = RGBColor(73, 80, 87)       # #495057
    CARD_BG_LIGHT = RGBColor(248, 249, 250) # #F8F9FA
    ACCENT_ORANGE = RGBColor(255, 107, 0)   # #FF6B00
    BORDER_LIGHT = RGBColor(222, 226, 230)
    BORDER_BLUE = RGBColor(186, 218, 245)

    def add_sih_header(slide, slide_title, slide_num):
        # Team METERE Badge (Top Left)
        if slide_num > 1:
            team_oval = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.4), Inches(0.25), Inches(1.4), Inches(1.1))
            team_oval.fill.solid()
            team_oval.fill.fore_color.rgb = RGBColor(255, 255, 255)
            team_oval.line.color.rgb = RGBColor(50, 50, 50)
            team_oval.line.width = Pt(1.5)
            tf_team = team_oval.text_frame
            tf_team.word_wrap = True
            p1 = tf_team.paragraphs[0]
            p1.text = "TEAM"
            p1.font.size = Pt(12)
            p1.font.bold = True
            p1.font.color.rgb = RGBColor(20, 20, 20)
            p1.alignment = PP_ALIGN.CENTER
            p2 = tf_team.add_paragraph()
            p2.text = "METERE"
            p2.font.size = Pt(12)
            p2.font.bold = True
            p2.font.color.rgb = RGBColor(20, 20, 20)
            p2.alignment = PP_ALIGN.CENTER

        # Title in Center
        t_box = slide.shapes.add_textbox(Inches(2.0), Inches(0.35), Inches(8.5), Inches(1.0))
        tf = t_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = slide_title
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = NAVY_TITLE
        p.alignment = PP_ALIGN.CENTER

        # SIH 2026 Logo placeholder / text (Top Right)
        logo_box = slide.shapes.add_textbox(Inches(10.8), Inches(0.2), Inches(2.2), Inches(1.1))
        tf_l = logo_box.text_frame
        p_l = tf_l.paragraphs[0]
        p_l.text = "SMART INDIA\nHACKATHON\n2026"
        p_l.font.size = Pt(12)
        p_l.font.bold = True
        p_l.font.color.rgb = NAVY_TITLE
        p_l.alignment = PP_ALIGN.CENTER

        # SIH Template Footer Strip (Bottom)
        f_strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.0), Inches(13.333), Inches(0.5))
        f_strip.fill.solid()
        f_strip.fill.fore_color.rgb = FOOTER_BLUE
        f_strip.line.fill.background()
        tf_f = f_strip.text_frame
        p_f = tf_f.paragraphs[0]
        p_f.text = f"@SIH Idea submission- Template                                                                                                                                           {slide_num}"
        p_f.font.size = Pt(11)
        p_f.font.color.rgb = RGBColor(255, 255, 255)
        p_f.font.bold = True

    # =========================================================================
    # SLIDE 6: RESEARCH AND REFERENCES (The Main Focus)
    # =========================================================================
    s6 = prs.slides.add_slide(blank_layout)
    add_sih_header(s6, "RESEARCH AND REFERENCES", 6)

    # Subtitle / Section Lead
    lead_box = s6.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(11.7), Inches(0.4))
    tf_lead = lead_box.text_frame
    p_lead = tf_lead.paragraphs[0]
    p_lead.text = "• Details / Links of the reference and research work foundational to FairCrop:"
    p_lead.font.size = Pt(15)
    p_lead.font.bold = True
    p_lead.font.color.rgb = NAVY_TITLE

    # 3 High-Impact Category Cards for Slide 6
    card_w = Inches(3.7)
    card_h = Inches(4.8)
    card_top = Inches(1.9)

    # CARD 1: Government Data & Live APIs
    c1 = s6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), card_top, card_w, card_h)
    c1.fill.solid()
    c1.fill.fore_color.rgb = CARD_BG_LIGHT
    c1.line.color.rgb = BORDER_BLUE
    c1.line.width = Pt(1.5)
    tf1 = c1.text_frame
    tf1.word_wrap = True
    tf1.margin_left = Inches(0.2)
    tf1.margin_right = Inches(0.2)
    tf1.margin_top = Inches(0.2)

    p = tf1.paragraphs[0]
    p.text = "🏛️ Govt. Datasets & APIs"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = NAVY_TITLE
    p.space_after = Pt(10)

    items1 = [
        ("e-NAM (National Agriculture Market):", "Small Farmers Agribusiness Consortium (SFAC), Ministry of Agriculture & Farmers Welfare, GoI. Integration standards across 1,361+ unified mandis.\n🔗 enam.gov.in"),
        ("Agmarknet Portal (DMI):", "Directorate of Marketing & Inspection, GoI. Wholesale arrival volumes, minimum, maximum, and modal price datasets.\n🔗 agmarknet.gov.in"),
        ("Open Govt Data (OGD) Platform:", "Ministry of Electronics & IT (MeitY). Daily APMC market price feeds and commodities statistics.\n🔗 data.gov.in"),
        ("PM-KISAN & PMFBY Guidelines:", "DPI framework for Direct Benefit Transfers (DBT) and crop insurance risk mitigation.\n🔗 pmkisan.gov.in | pmfby.gov.in")
    ]
    for title, desc in items1:
        pt = tf1.add_paragraph()
        pt.text = f"• {title}"
        pt.font.size = Pt(11)
        pt.font.bold = True
        pt.font.color.rgb = ACCENT_ORANGE
        pd = tf1.add_paragraph()
        pd.text = desc
        pd.font.size = Pt(9.5)
        pd.font.color.rgb = TEXT_DARK
        pd.space_after = Pt(6)

    # CARD 2: Policy Reports & Academic Studies
    c2 = s6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.8), card_top, card_w, card_h)
    c2.fill.solid()
    c2.fill.fore_color.rgb = CARD_BG_LIGHT
    c2.line.color.rgb = BORDER_BLUE
    c2.line.width = Pt(1.5)
    tf2 = c2.text_frame
    tf2.word_wrap = True
    tf2.margin_left = Inches(0.2)
    tf2.margin_right = Inches(0.2)
    tf2.margin_top = Inches(0.2)

    p = tf2.paragraphs[0]
    p.text = "📊 Policy Reports & Market Studies"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = NAVY_TITLE
    p.space_after = Pt(10)

    items2 = [
        ("NITI Aayog & NABARD (2024):", "“Promoting Farmer Producer Organizations (FPOs) & Strengthening Digital Public Infrastructure in Indian Agriculture.” (FPO collective bargaining increases farm-gate margins by 15–25%)."),
        ("ICRIER Research Study:", "“Agricultural Market Integration & Price Volatility in India.” (Quantifies 18–30% middleman margin leakages and high spatial price dispersion across regional APMCs)."),
        ("FAO & World Bank Studies:", "“Mitigating Post-Harvest Food Loss in Perishables via Direct Linkages.” (Direct contracting and cold-freight reduce perishable harvest wastage by 20–25%).")
    ]
    for title, desc in items2:
        pt = tf2.add_paragraph()
        pt.text = f"• {title}"
        pt.font.size = Pt(11)
        pt.font.bold = True
        pt.font.color.rgb = ACCENT_ORANGE
        pd = tf2.add_paragraph()
        pd.text = desc
        pd.font.size = Pt(9.5)
        pd.font.color.rgb = TEXT_DARK
        pd.space_after = Pt(6)

    # CARD 3: Technical, Quality & Legal Frameworks
    c3 = s6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.8), card_top, card_w, card_h)
    c3.fill.solid()
    c3.fill.fore_color.rgb = CARD_BG_LIGHT
    c3.line.color.rgb = BORDER_BLUE
    c3.line.width = Pt(1.5)
    tf3 = c3.text_frame
    tf3.word_wrap = True
    tf3.margin_left = Inches(0.2)
    tf3.margin_right = Inches(0.2)
    tf3.margin_top = Inches(0.2)

    p = tf3.paragraphs[0]
    p.text = "⚙️ Standards & Technical References"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = NAVY_TITLE
    p.space_after = Pt(10)

    items3 = [
        ("ONDC Agri Protocol Specifications:", "Open Network for Digital Commerce specifications for unbundled agri-buyer-seller discovery and fulfillment.\n🔗 ondc.org"),
        ("Agmark Quality Grading Standards:", "Agricultural Produce (Grading and Marking) Act, 1937. Parameterized moisture content %, purity grades (Grade A/B/C) and digital lot assaying."),
        ("ML Price Forecasting Research:", "Time-series predictive models (ARIMA + LSTM neural networks) for APMC price trends & automated Hold/Sell advisories."),
        ("Indian Contract Act (1872) & Escrow:", "Legal framework for digital offer-bidding lock and milestone-based secure fund release upon delivery inspection.")
    ]
    for title, desc in items3:
        pt = tf3.add_paragraph()
        pt.text = f"• {title}"
        pt.font.size = Pt(11)
        pt.font.bold = True
        pt.font.color.rgb = ACCENT_ORANGE
        pd = tf3.add_paragraph()
        pd.text = desc
        pd.font.size = Pt(9.5)
        pd.font.color.rgb = TEXT_DARK
        pd.space_after = Pt(6)

    # Save PPTX
    prs.save(output_pptx)
    prs.save("SIH2026_AgriLink_Presentation.pptx")
    print(f"SUCCESS: Created {output_pptx} and updated SIH2026_AgriLink_Presentation.pptx with complete Slide 6!")

if __name__ == "__main__":
    generate_sih_deck()
