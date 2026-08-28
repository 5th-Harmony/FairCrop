# -*- coding: utf-8 -*-
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def build_viva_deck(output_path="SIH2026_Technical_Viva_Guide.pptx"):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Colors
    NAVY = RGBColor(16, 44, 87)       # #102C57
    EMERALD = RGBColor(15, 81, 50)    # #0F5132
    ORANGE = RGBColor(255, 107, 0)    # #FF6B00
    LIGHT_BG = RGBColor(248, 250, 252)# #F8FAFC
    DARK_BG = RGBColor(15, 23, 42)    # #0F172A
    CARD_BG = RGBColor(255, 255, 255)
    BORDER_CLR = RGBColor(226, 232, 240)
    TEXT_DARK = RGBColor(30, 41, 59)
    TEXT_MUTED = RGBColor(100, 116, 139)

    def set_bg(slide, color):
        s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        s.fill.solid()
        s.fill.fore_color.rgb = color
        s.line.fill.background()
        return s

    def add_header(slide, title, category="SIH 2026 • TECHNICAL VIVA & JUDGES Q&A CHEATSHEET"):
        c_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.35))
        tf_c = c_box.text_frame
        p_c = tf_c.paragraphs[0]
        p_c.text = category.upper()
        p_c.font.size = Pt(11)
        p_c.font.bold = True
        p_c.font.color.rgb = ORANGE

        t_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.7))
        tf_t = t_box.text_frame
        p_t = tf_t.paragraphs[0]
        p_t.text = title
        p_t.font.size = Pt(24)
        p_t.font.bold = True
        p_t.font.color.rgb = NAVY

    def add_card(slide, left, top, width, height, title, bullet_list, accent=BORDER_CLR, bg=CARD_BG, title_color=NAVY):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        card.fill.solid()
        card.fill.fore_color.rgb = bg
        card.line.color.rgb = accent
        card.line.width = Pt(1.5)

        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.25)
        tf.margin_right = Inches(0.25)
        tf.margin_top = Inches(0.22)
        tf.margin_bottom = Inches(0.22)

        p0 = tf.paragraphs[0]
        p0.text = title
        p0.font.size = Pt(15)
        p0.font.bold = True
        p0.font.color.rgb = title_color
        p0.space_after = Pt(8)

        for item in bullet_list:
            p = tf.add_paragraph()
            p.text = item
            p.font.size = Pt(11)
            p.font.color.rgb = TEXT_DARK
            p.space_after = Pt(5)

    # ──────────────────────────────────────────────────────────
    # SLIDE 1: Title Slide (High-Confidence Viva Prep)
    # ──────────────────────────────────────────────────────────
    s1 = prs.slides.add_slide(blank_layout)
    set_bg(s1, DARK_BG)

    tb = s1.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(11.3), Inches(4.5))
    tf1 = tb.text_frame
    tf1.word_wrap = True
    
    p = tf1.paragraphs[0]
    p.text = "SMART INDIA HACKATHON 2026"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = ORANGE
    p.space_after = Pt(12)

    p2 = tf1.add_paragraph()
    p2.text = "FairCrop — Technical Viva & Judge Q&A Guide"
    p2.font.size = Pt(32)
    p2.font.bold = True
    p2.font.color.rgb = RGBColor(255, 255, 255)
    p2.space_after = Pt(16)

    p3 = tf1.add_paragraph()
    p3.text = "Everything you need to know to answer all technical, architectural, and domain questions with full confidence."
    p3.font.size = Pt(15)
    p3.font.color.rgb = RGBColor(203, 213, 225)
    p3.space_after = Pt(24)

    p4 = tf1.add_paragraph()
    p4.text = "Team: METERE  •  Problem Statement: Strengthening market linkages and price discovery for farmers"
    p4.font.size = Pt(13)
    p4.font.bold = True
    p4.font.color.rgb = RGBColor(52, 211, 153)

    # ──────────────────────────────────────────────────────────
    # SLIDE 2: 30-Second Elevator Pitch & Core Problem
    # ──────────────────────────────────────────────────────────
    s2 = prs.slides.add_slide(blank_layout)
    set_bg(s2, LIGHT_BG)
    add_header(s2, "1. The 30-Second Pitch & What Problem It Solves")

    add_card(s2, 0.8, 1.6, 5.6, 5.2, "🗣️ What is FairCrop? (Say This in 2 Lines)", [
        "“FairCrop is a smart digital agricultural marketplace that connects farmers directly with verified buyers.”",
        "“It eliminates middlemen, provides AI-driven price forecasts (Hold vs Sell), and guarantees 100% safe payments through a secure escrow system.”",
        "Key Mission: Increase farmer income by 15–25% and reduce post-harvest perishable wastage by 20%."
    ], accent=ORANGE)

    add_card(s2, 6.8, 1.6, 5.7, 5.2, "⚠️ The 3 Core Problems We Solve", [
        "1. Middleman Exploitation: Traditional commission agents charge 15–30% hidden margins from farmers.",
        "2. Lack of Price Visibility: Farmers don't know prices in neighboring mandis and sell at distress rates.",
        "3. Payment Default Risk: Buyers delay or default on payments after taking produce; FairCrop locks payment in Escrow upfront."
    ], accent=NAVY)

    # ──────────────────────────────────────────────────────────
    # SLIDE 3: The Complete Tech Stack (Explained in Plain English)
    # ──────────────────────────────────────────────────────────
    s3 = prs.slides.add_slide(blank_layout)
    set_bg(s3, LIGHT_BG)
    add_header(s3, "2. Complete Technology Stack (Plain English Breakdown)")

    add_card(s3, 0.8, 1.6, 3.6, 5.2, "🌐 Frontend (User Interface)", [
        "• HTML5 & CSS3: Clean glassmorphism design, ultra-fast loading for rural 3G/4G.",
        "• Vanilla JavaScript: Fast, zero dependency bloat, responsive interactivity.",
        "• Next.js & React: Modern web app portal for institutional buyer dashboards.",
        "• React Native / Expo: Mobile app for farmers on Android smartphones."
    ], accent=RGBColor(59, 130, 246))

    add_card(s3, 4.8, 1.6, 3.6, 5.2, "⚙️ Backend (The Engine)", [
        "• Python FastAPI: Asynchronous, ultra-fast backend framework (exceeds Flask/Django in speed).",
        "• SQLAlchemy 2.0 ORM: Clean database models & relationship mapping.",
        "• Pydantic v2: Strict data validation and type safety.",
        "• JWT & Bcrypt: Secure token login and encrypted password storage."
    ], accent=EMERALD)

    add_card(s3, 8.8, 1.6, 3.6, 5.2, "🤖 Database, AI & Data", [
        "• Database: SQLite for development; PostgreSQL for scalable cloud production.",
        "• AI/ML: Time-series predictive models (ARIMA + LSTM) for 7-day crop price trends.",
        "• Data Feeds: e-NAM (1,361 mandis) & Agmarknet government open APIs.",
        "• Localization: 22 official Indian languages supported."
    ], accent=ORANGE)

    # ──────────────────────────────────────────────────────────
    # SLIDE 4: How It Works — The 7-Stage Transaction Flow
    # ──────────────────────────────────────────────────────────
    s4 = prs.slides.add_slide(blank_layout)
    set_bg(s4, LIGHT_BG)
    add_header(s4, "3. End-to-End Workflow & 7-Stage Escrow Flow")

    add_card(s4, 0.8, 1.6, 5.6, 5.2, "🌾 Step-by-Step Platform Journey", [
        "1. Registration: Farmer/FPO signs up via mobile OTP & Aadhaar.",
        "2. Lot Creation: Farmer lists crop, quantity (quintals), quality grade (Agmark A/B/C), and expected price.",
        "3. AI Matchmaking: System automatically matches the lot with nearby verified institutional buyers.",
        "4. Bidding & Counter-Offers: Real-time negotiation until a fair price is locked."
    ], accent=EMERALD)

    add_card(s4, 6.8, 1.6, 5.7, 5.2, "🔒 The 7-Stage Escrow Security Flow", [
        "1. Bid Accepted → Contract digitally locked.",
        "2. Escrow Deposited → Buyer locks 100% payment upfront.",
        "3. Dispatched → Goods shipped with GPS tracking.",
        "4. Delivered → Produce arrives at buyer warehouse.",
        "5. Quality Verified → Funds instantly released to farmer.",
        "6. Dispute / Pause → If moisture/grade mismatch, admin mediates.",
        "7. Revert / Cancel → Funds refunded; lot reverts to available."
    ], accent=ORANGE)

    # ──────────────────────────────────────────────────────────
    # SLIDE 5: Top 6 Technical Viva Questions & Winning Answers
    # ──────────────────────────────────────────────────────────
    s5 = prs.slides.add_slide(blank_layout)
    set_bg(s5, LIGHT_BG)
    add_header(s5, "4. Most Common Technical Questions Judges Ask (Part 1)")

    add_card(s5, 0.8, 1.6, 5.6, 2.5, "Q1: Why FastAPI instead of Django/Flask?", [
        "Answer: FastAPI is asynchronous (high concurrency), up to 3x faster than Flask, auto-generates interactive Swagger API docs, and has native Pydantic validation."
    ], accent=RGBColor(59, 130, 246))

    add_card(s5, 6.8, 1.6, 5.7, 2.5, "Q2: Where do you get live mandi prices?", [
        "Answer: We integrate with official Government APIs — e-NAM (National Agriculture Market) and Agmarknet (Directorate of Marketing & Inspection) daily feeds."
    ], accent=EMERALD)

    add_card(s5, 0.8, 4.3, 5.6, 2.5, "Q3: How does the AI Price Forecast work?", [
        "Answer: It uses historical APMC time-series data to predict 7-day future price bands and provides actionable 'HOLD or SELL' advisories to farmers."
    ], accent=ORANGE)

    add_card(s5, 6.8, 4.3, 5.7, 2.5, "Q4: How do you prevent fraud/cheating?", [
        "Answer: Our 7-stage Escrow system holds 100% buyer money in escrow before shipping. Farmers are guaranteed payment upon verified delivery."
    ], accent=NAVY)

    # ──────────────────────────────────────────────────────────
    # SLIDE 6: Common Viva Questions (Part 2 — Rural Scale & Business)
    # ──────────────────────────────────────────────────────────
    s6 = prs.slides.add_slide(blank_layout)
    set_bg(s6, LIGHT_BG)
    add_header(s6, "5. Scalability, Rural Connectivity & Business Model")

    add_card(s6, 0.8, 1.6, 5.6, 2.5, "Q5: How will uneducated rural farmers use this?", [
        "Answer: 22 official Indian regional languages with native scripts, voice search capability, and integration with local FPOs (Farmer Producer Organizations) for bulk assisted selling."
    ], accent=EMERALD)

    add_card(s6, 6.8, 1.6, 5.7, 2.5, "Q6: What if internet connectivity is poor?", [
        "Answer: Lightweight client architecture with local storage caching, fast-loading static assets, and SMS/IVR voice integration in our scale roadmap."
    ], accent=RGBColor(59, 130, 246))

    add_card(s6, 0.8, 4.3, 5.6, 2.5, "Q7: What is the Business / Revenue Model?", [
        "Answer: 100% FREE for smallholder farmers. A small 0.5%–1.5% convenience fee charged to institutional bulk buyers for verified quality sourcing and logistics."
    ], accent=ORANGE)

    add_card(s6, 6.8, 4.3, 5.7, 2.5, "Q8: How does this scale nationally?", [
        "Answer: Designed to plug into ONDC (Open Network for Digital Commerce) agri-protocol and India's Digital Public Infrastructure (DPI) for agriculture."
    ], accent=NAVY)

    prs.save(output_path)
    print(f"SUCCESS: Generated {output_path}")

if __name__ == "__main__":
    build_viva_deck()
